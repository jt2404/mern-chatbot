# Importing libraries
import os
from langchain.vectorstores import FAISS
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import docx2txt
from dotenv import load_dotenv
load_dotenv()
secret_key = os.environ.get('SECRET_KEY')

# RAG model using langchain


def read_pdf_file(file_path):
    with open(file_path, 'rb') as file:
        pdf_reader = PdfReader(file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text() + '\n'
        return text.strip()


def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_ppt(ppt_path):
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    text =  '\n'.join(full_text)
    return text


def read_tex_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


def extract_text_from_doc(file_path):
  text = docx2txt.process(file_path)
  return text


def get_file_type(file_name):
    _, file_extension = os.path.splitext(file_name)
    return file_extension.lower()


def extract_text_from_file(file_name):
    file_type = get_file_type(file_name)
    if file_type == '.pdf':
        return read_pdf_file(file_name)
    elif file_type == '.pptx':
        return extract_text_from_pptx(file_name)
    elif file_type == '.ppt':
        return extract_text_from_ppt(file_name)
    elif file_type == '.docx':
        return extract_text_from_docx(file_name)
    elif file_type == '.doc':
        return extract_text_from_doc(file_name)
    elif file_type == '.tex':
        return read_tex_file(file_name)
    else:
        return "Unsupported file type."


def text_to_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(
              chunk_size=512,
              chunk_overlap=32,
              length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def processing_file(chunks,file_name):
    print(secret_key)
    os.environ["OPENAI_API_KEY"] = secret_key
    embeddings = OpenAIEmbeddings(openai_api_key=secret_key)
    docsearch = FAISS.from_texts(chunks, embeddings)
    chain = load_qa_chain(OpenAI(), chain_type="stuff")
    print(f"\nProcessing your file: {file_name}\n")
    return docsearch, chain



def question_answering2(docsearch,chain):
  while True:
    question = "Please Give the summary of given pdf"
    docs = docsearch.similarity_search(question)
    ans = chain.run(input_documents=docs, question=question)
    return ans
global answer
def question_answering(docsearch,chain,que):
  while True:
    # question = input("Enter your question (or type 'quit' to exit): ")
    question = que
    if question.lower() == 'quit':
        break
    docs = docsearch.similarity_search(question)
    answer = chain.run(input_documents=docs, question=question)
    # print(f"Answer: {answer}\n")
    return answer
    # Store the question-answer pair
    # answers.append({"question": question, "answer": answer})

    # # Store the answers for the current LaTeX file
    # all_answers[file_path] = answers

    # return all_answers

def show_ans(answer):
    print("in show")
    ans = answer
    return ans


